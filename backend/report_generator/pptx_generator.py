"""PowerPoint Generator - creates PPTX reports."""
from __future__ import annotations

import json
import logging
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config.models import CompanyConfig
from config.settings import settings

logger = logging.getLogger(__name__)

# Colors
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x33, 0x33, 0x33)
C_GRAY_ALT = RGBColor(0xF2, 0xF2, 0xF2)
C_GRAY_LINE = RGBColor(0xCC, 0xCC, 0xCC)
C_POS = RGBColor(0x1A, 0x7A, 0x1A)
C_NEG = RGBColor(0xCC, 0x00, 0x00)
C_RED_LIGHT = RGBColor(0xFF, 0xEB, 0xEB)
C_TBL_HDR = RGBColor(0x44, 0x44, 0x44)

# Fonts
FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_BULLET = "Commissioner"

# Sizes
SZ_HEADER = 18
SZ_COVER_T = 22
SZ_COVER_S = 10
SZ_TBL_HDR = 9
SZ_TBL_CELL = 8
SZ_BULLET = 12.5
SZ_PARA = 11

# Layout
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
CONTENT_TOP = 0.82
CONTENT_BOT = 5.50
CONTENT_H = CONTENT_BOT - CONTENT_TOP
SLIDE_BOTTOM = 5.625
TABLE_LEFT = 1.0
TABLE_RIGHT = 9.0
TABLE_W = TABLE_RIGHT - TABLE_LEFT


class PPTXGenerator:
    """Generates PowerPoint SEO reports."""

    def __init__(
        self,
        company: CompanyConfig,
        charts_dir: Path | None = None,
    ) -> None:
        self.company = company
        self.charts_dir = charts_dir or settings.charts_dir

        # Convert hex colors
        self.bar_color = self._hex_to_rgb(company.header_color)
        self.accent_color = self._hex_to_rgb(company.accent_color)

    @classmethod
    def for_company(cls, company: CompanyConfig, company_key: str) -> "PPTXGenerator":
        """Construct with correctly scoped charts directory per company key."""
        charts_dir = settings.charts_dir / company_key.lower()
        return cls(company, charts_dir)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _px(self, v: float) -> Inches:
        return Inches(v)

    def _add_slide(self, prs: Presentation):
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _add_rect(self, slide, l, t, w, h, fill=None, line=None, lw=0.25):
        s = slide.shapes.add_shape(1, self._px(l), self._px(t), self._px(w), self._px(h))
        if fill:
            s.fill.solid()
            s.fill.fore_color.rgb = fill
        else:
            s.fill.background()
        if line:
            s.line.color.rgb = line
            s.line.width = Pt(lw)
        else:
            s.line.fill.background()
        return s

    def _add_text(self, slide, l, t, w, h, text, size=11, bold=False, italic=False,
                  color=C_DARK, align=PP_ALIGN.LEFT, font=FONT_BODY, wrap=True):
        tb = slide.shapes.add_textbox(self._px(l), self._px(t), self._px(w), self._px(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
        return tb

    def _add_image(self, slide, path, l, t, w, h):
        if os.path.exists(path):
            slide.shapes.add_picture(path, self._px(l), self._px(t), self._px(w), self._px(h))
        else:
            self._add_rect(slide, l, t, w, h, fill=C_GRAY_ALT)
            self._add_text(slide, l+0.1, t+h/2-0.15, w-0.2, 0.3,
                          f"[Missing chart: {os.path.basename(path)}]",
                          size=8, align=PP_ALIGN.CENTER)

    def _short_m(self, label: str) -> str:
        parts = label.split()
        return f"{parts[0][:3]} {parts[1][2:]}" if len(parts) >= 2 else label

    def _clean_text(self, text: str) -> str:
        text = text.replace("—", ",").replace("–", ",")
        text = re.sub(r"^\s*[-—–]\s*", "", text)
        return text.strip()

    def _split_bullets(self, text: str) -> List[str]:
        return [self._clean_text(l) for l in text.split("\n") if l.strip()]

    def _header_bar(self, slide, title):
        self._add_rect(slide, 0, 0, 10, 0.75, fill=self.bar_color)
        self._add_text(slide, 0.25, 0.10, 9.5, 0.57,
                       title, size=SZ_HEADER, bold=True,
                       color=C_WHITE, font=FONT_TITLE)

    def _add_native_table(self, slide, t, col_widths, headers, rows, delta_col=None, rh=0.30, cell_rh=0.26, row_fill=None):
        num_cols = len(headers)
        num_rows = len(rows) + 1
        
        table_shape = slide.shapes.add_table(
            num_rows, num_cols, self._px(TABLE_LEFT), self._px(t), 
            self._px(sum(col_widths)), self._px(rh + cell_rh * len(rows))
        )
        table = table_shape.table

        # Set column widths
        for i, w in enumerate(col_widths):
            table.columns[i].width = self._px(w)

        # Header Row
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_TBL_HDR
            
            cell.margin_left = self._px(0.05)
            cell.margin_right = self._px(0.05)
            cell.margin_top = self._px(0.02)
            cell.margin_bottom = self._px(0.02)
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(SZ_TBL_HDR)
                    run.font.bold = True
                    run.font.color.rgb = C_WHITE
                    run.font.name = FONT_BODY

        # Data Rows
        for r_idx, row in enumerate(rows):
            fill_color = row_fill if row_fill is not None else (C_GRAY_ALT if r_idx % 2 == 1 else C_WHITE)
            for c_idx, cell_val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(cell_val)
                
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_color
                
                cell.margin_left = self._px(0.05)
                cell.margin_right = self._px(0.05)
                cell.margin_top = self._px(0.02)
                cell.margin_bottom = self._px(0.02)
                
                txt_col = C_DARK
                if delta_col is not None and c_idx == delta_col:
                    try:
                        v = float(str(cell_val).replace("+", "").replace(",", "").replace("%", ""))
                        txt_col = C_POS if v >= 0 else C_NEG
                    except ValueError:
                        pass
                        
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(SZ_TBL_CELL)
                        run.font.color.rgb = txt_col
                        run.font.name = FONT_BODY
                        
        return t + rh + (cell_rh * len(rows))

    def _bullet_textbox(self, slide, l, t, w, h, lines, size=SZ_BULLET):
        tb = slide.shapes.add_textbox(self._px(l), self._px(t), self._px(w), self._px(h))
        tf = tb.text_frame
        tf.word_wrap = True

        for i, raw_line in enumerate(lines):
            line = self._clean_text(raw_line)
            if not line:
                continue
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(6)

            bullet_run = p.add_run()
            bullet_run.text = "•  "
            bullet_run.font.size = Pt(size)
            bullet_run.font.name = FONT_BULLET
            bullet_run.font.color.rgb = C_DARK
            bullet_run.font.bold = True

            text_run = p.add_run()
            text_run.text = line
            text_run.font.size = Pt(size)
            text_run.font.name = FONT_BULLET
            text_run.font.color.rgb = C_DARK

        return tb

    def generate(
        self,
        data: Dict[str, Any],
        texts: Dict[str, str],
        output_path: Path | None = None,
        slide_list: List[str] | None = None,
    ) -> Path:
        """Generate the complete PPTX report."""
        logger.info("Building PPTX for %s", self.company.name)

        # Determine output path
        if output_path is None:
            months = list(data["traffic"]["monthly"].keys())
            last_month = months[-1] if months else "Unknown"
            month_slug = last_month.replace(" ", "_")
            output_path = settings.reports_dir / self.company.name.replace(" ", "_") / f"SEO_Report_{self.company.name.replace(' ', '_')}_{month_slug}.pptx"
        else:
            output_path = Path(output_path)

        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Create presentation
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        # Fix 1: Safe month_label derivation — always from actual data keys
        monthly_keys = list(data["traffic"]["monthly"].keys())
        month_label = monthly_keys[-1] if monthly_keys else "Unknown"

        # Fix 2: Explicit argument binding helper — prevents closure capture issues
        # when slides are removed or reordered.
        def _call(fn, *args, **kwargs):
            return lambda: fn(*args, **kwargs)

        # Map slide IDs to explicitly-bound callables
        slide_mapping = {
            "cover":                       _call(self._slide_cover, prs, data, month_label),
            "traffic_overview":            _call(self._slide_traffic, prs, data, month_label),
            "6month_summary":              _call(self._slide_6month_summary, prs, texts, month_label),
            "monthly_summary":             _call(self._slide_monthly_summary, prs, texts, month_label),
            "content_trending_up":         _call(self._slide_content_trending_up, prs, data, texts, month_label),
            "content_trending_down":       _call(self._slide_content_trending_down, prs, data, texts, month_label),
            "top_queries":                 _call(self._slide_top_queries, prs, data, month_label),
            "keyword_table_branded":       _call(self._slide_keyword_table, prs, data, "branded", month_label),
            "keyword_analysis_branded":    _call(self._slide_keyword_analysis, prs, texts, "branded", month_label),
            "keyword_table_nonbranded":    _call(self._slide_keyword_table, prs, data, "non_branded", month_label),
            "keyword_analysis_nonbranded": _call(self._slide_keyword_analysis, prs, texts, "non_branded", month_label),
            "branded_vs_nonbranded":       _call(self._slide_branded_vs_nonbranded, prs, data, month_label),
            "pages_table_nonblog":         _call(self._slide_pages_table, prs, data, "non_blog_pages", month_label),
            "pages_analysis_nonblog":      _call(self._slide_pages_analysis, prs, data, "non_blog_pages", month_label),
            "pages_table_blog":            _call(self._slide_pages_table, prs, data, "blog_pages", month_label),
            "pages_analysis_blog":         _call(self._slide_pages_analysis, prs, data, "blog_pages", month_label),
            "channels":                    _call(self._slide_channels, prs, data, texts, month_label),
            "channels_charts":             _call(self._slide_channels_charts, prs, month_label),
            "core_web_vitals":             _call(self._slide_core_web_vitals, prs, data, month_label),
        }

        # Build slides in the requested order (or full default order if none given)
        if slide_list is None:
            slide_list = list(slide_mapping.keys())

        for slide_id in slide_list:
            if slide_id in slide_mapping:
                slide_mapping[slide_id]()
            else:
                logger.warning("Unknown slide ID requested: %s", slide_id)

        # Save
        prs.save(str(output_path))
        logger.info("PPTX saved: %s", output_path)
        return output_path

    def _slide_cover(self, prs, data, month_label):
        slide = self._add_slide(prs)
        self._add_rect(slide, 0, 0, 10, 0.75, fill=self.bar_color)
        self._add_rect(slide, 3.4, 1.45, 3.3, 1.6, fill=C_DARK)
        self._add_text(slide, 3.45, 2.0, 3.2, 0.7,
                       self.company.name, size=20, bold=True,
                       color=C_WHITE, align=PP_ALIGN.CENTER)
        self._add_text(slide, 0.4, 3.3, 9.2, 0.75,
                       f"SEO Performance Report: {month_label}",
                       size=SZ_COVER_T, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        self._add_text(slide, 0.4, 4.2, 9.2, 0.35,
                       f"Generated: {datetime.now().strftime('%B %d, %Y')}  |  {self.company.gsc_url}",
                       size=SZ_COVER_S, color=self.accent_color, align=PP_ALIGN.CENTER)

    def _slide_traffic(self, prs, data, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, "Traffic Overview (Past 6 Months)")
        monthly = data["traffic"]["monthly"]

        cw = [1.3, 0.9, 1.0]
        rows = [[self._short_m(m), f"{monthly[m]['clicks']:,}", f"{monthly[m]['impressions']:,}"]
                for m in monthly.keys()]
        self._add_native_table(slide, CONTENT_TOP, cw, ["Month", "Clicks", "Impressions"], rows)

        chart_path = str(self.charts_dir / "chart_clicks_only.png")
        self._add_image(slide, chart_path, TABLE_LEFT + 3.3, CONTENT_TOP, TABLE_W - 3.4, CONTENT_H)

    def _slide_6month_summary(self, prs, texts, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, "SEO Performance: Past 6 Months")
        paras = self._split_bullets(texts.get("six_month_summary", ""))
        self._bullet_textbox(slide, TABLE_LEFT, CONTENT_TOP, TABLE_W, CONTENT_H, paras)

    def _slide_monthly_summary(self, prs, texts, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, f"SEO Performance: {month_label}")
        paras = self._split_bullets(texts.get("monthly_summary", ""))
        self._bullet_textbox(slide, TABLE_LEFT, CONTENT_TOP, TABLE_W, CONTENT_H, paras)

    def _slide_content_trending_up(self, prs, data, texts, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, "Content Trending Up: Past 28 Days")
        pages = data.get("page_gainers", [])[:5]

        cw = [4.5, 1.0, 1.0, 1.5]

        rows = []
        for g in pages:
            url = g["page"].replace("https://www.", "").replace("https://", "")[:70]
            ctype = "Blog Post" if "/blog/" in g["page"] else "Service Page"
            if "/careers" in g["page"]:
                ctype = "Careers"
            rows.append([url, str(g["current_clicks"]),
                         f"+{g['click_delta']}" if g["click_delta"] >= 0 else str(g["click_delta"]),
                         ctype])

        bottom = self._add_native_table(slide, CONTENT_TOP, cw, ["Page URL", "Clicks", "MoM Change", "Content Type"], rows, delta_col=2)
        buls = self._split_bullets(texts.get("content_trending_up", ""))
        self._bullet_textbox(slide, TABLE_LEFT, bottom + 0.12, TABLE_W,
                             SLIDE_BOTTOM - bottom - 0.15, buls)

    def _slide_content_trending_down(self, prs, data, texts, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, "Content Trending Down: Past 28 Days")
        pages = data.get("page_decliners", [])[:5]

        cw = [4.5, 1.0, 1.0, 1.5]

        rows = []
        for g in pages:
            url = g["page"].replace("https://www.", "").replace("https://", "")[:70]
            ctype = "Blog Post" if "/blog/" in g["page"] else "Service Page"
            rows.append([url, str(g["current_clicks"]),
                         f"+{g['click_delta']}" if g["click_delta"] >= 0 else str(g["click_delta"]),
                         ctype])

        bottom = self._add_native_table(slide, CONTENT_TOP, cw, ["Page URL", "Clicks", "MoM Change", "Content Type"], rows, delta_col=2, row_fill=C_RED_LIGHT)
        buls = self._split_bullets(texts.get("content_trending_down", ""))
        self._bullet_textbox(slide, TABLE_LEFT, bottom + 0.12, TABLE_W,
                             SLIDE_BOTTOM - bottom - 0.15, buls)

    def _slide_top_queries(self, prs, data, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, "Top Queries: Period Overview")
        kws = data.get("keywords", {}).get("branded", [])[:5] + \
              data.get("keywords", {}).get("non_branded", [])[:5]

        cw = [3.2, 1.0, 1.4, 1.0, 1.4]
        rows = [[k["query"][:60], str(k["clicks"]), f"{k['impressions']:,}",
                 f"{k['ctr']*100:.1f}%", f"{k['position']:.1f}"] for k in kws]
        self._add_native_table(slide, CONTENT_TOP, cw, ["Query", "Clicks", "Impressions", "CTR", "Avg Position"], rows)

    def _slide_keyword_table(self, prs, data, kw_type, month_label):
        slide = self._add_slide(prs)
        label = "Branded" if kw_type == "branded" else "Non-Branded"
        self._header_bar(slide, f"Top Performing {label} Keywords: {month_label}")
        kws = data.get("keywords", {}).get(kw_type, [])[:10]

        cw = [3.2, 1.0, 1.4, 1.0, 1.4]
        rows = [[k["query"][:60], str(k["clicks"]), f"{k['impressions']:,}",
                 f"{k['ctr']*100:.1f}%", f"{k['position']:.1f}"] for k in kws]
        self._add_native_table(slide, CONTENT_TOP, cw, ["Query", "Clicks", "Impressions", "CTR", "Avg Position"], rows)

    def _slide_keyword_analysis(self, prs, texts, kw_type, month_label):
        slide = self._add_slide(prs)
        label = "Branded" if kw_type == "branded" else "Non-Branded"
        self._header_bar(slide, f"Top Performing {label} Keywords: {month_label}")
        key = "branded_kw_analysis" if kw_type == "branded" else "nonbranded_kw_analysis"
        buls = self._split_bullets(texts.get(key, ""))
        self._bullet_textbox(slide, TABLE_LEFT, CONTENT_TOP, TABLE_W, CONTENT_H, buls)

    def _slide_branded_vs_nonbranded(self, prs, data, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, "Branded vs Non-Branded Keyword Traffic")
        bvnb = data["traffic"]["branded_vs_nonbranded"]

        cw = [1.4, 0.9, 0.9]
        rows = [[self._short_m(m), str(bvnb[m]["branded"]), str(bvnb[m]["non_branded"])]
                for m in bvnb.keys()]
        self._add_native_table(slide, CONTENT_TOP, cw, ["Month", "Branded", "Non-Branded"], rows)

        chart_path = str(self.charts_dir / "chart_branded_nonbranded.png")
        self._add_image(slide, chart_path, TABLE_LEFT + 3.3, CONTENT_TOP, TABLE_W - 3.4, CONTENT_H)

    def _slide_pages_table(self, prs, data, page_type, month_label):
        slide = self._add_slide(prs)
        label = "Blog Pages" if page_type == "blog_pages" else "Pages"
        self._header_bar(slide, f"Top Performing {label}: {month_label}")
        pages = data.get("pages", {}).get(page_type, [])[:10]

        cw = [3.0, 1.0, 1.4, 1.0, 1.6]
        rows = [[p.get("title", p["url"])[:55], str(p["clicks"]), f"{p['impressions']:,}",
                 f"{p['ctr']*100:.1f}%", f"{p['position']:.1f}"] for p in pages]
        self._add_native_table(slide, CONTENT_TOP, cw, ["Page", "Clicks", "Impressions", "CTR", "Avg Position"], rows)

    def _slide_pages_analysis(self, prs, data, page_type, month_label):
        slide = self._add_slide(prs)
        label = "Blog Pages" if page_type == "blog_pages" else "Pages"
        self._header_bar(slide, f"Top Performing {label}: {month_label}")
        pages = data.get("pages", {}).get(page_type, [])[:5]

        # Fix 3: Guard against empty data — render a clear placeholder instead of blank slide
        if not pages:
            self._add_text(
                slide, TABLE_LEFT, CONTENT_TOP, TABLE_W, CONTENT_H,
                f"No {label.lower()} data available for this reporting period.",
                size=11, color=C_DARK
            )
            return

        buls = []
        if pages:
            top = pages[0]
            title = top.get("title", top["url"])
            buls.append(
                f"Top page: {title} recorded {top['clicks']:,} clicks with a {top['ctr']*100:.2f}% CTR at an average position of {top['position']:.1f}. This page drives the majority of organic traffic for this category."
            )
            for p in pages[1:4]:
                buls.append(
                    f"{p.get('title', p['url'])[:60]} delivered {p['clicks']} clicks at {p['ctr']*100:.1f}% CTR from position {p['position']:.1f}."
                )
            buls.append(
                "Pages with more than 3,000 impressions and sub-1% CTR represent the strongest opportunity for title tag and meta description optimisation to convert existing visibility into clicks."
            )

        self._bullet_textbox(slide, TABLE_LEFT, CONTENT_TOP, TABLE_W, CONTENT_H, buls)

    def _slide_channels(self, prs, data, texts, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, "Organic Channel Comparison")

        ch_data = data.get("channels", [])
        if not ch_data:
            self._add_text(slide, TABLE_LEFT, CONTENT_TOP, TABLE_W, CONTENT_H, "No channel data available")
            return

        ch_keys = ["Google", "ChatGPT", "Bing", "LinkedIn", "Perplexity", "Direct", "Email", "Claude"]
        months = [c["month"] for c in ch_data]

        name_w = 1.4
        mon_w = round((TABLE_W - name_w) / len(months), 4)
        cw = [name_w] + [mon_w] * len(months)
        cw[-1] = round(TABLE_W - sum(cw[:-1]), 4)

        rows = [[key] + [str(c.get(key, 0)) for c in ch_data] for key in ch_keys]
        bottom = self._add_native_table(slide, CONTENT_TOP, cw, ["Channel"] + months, rows, rh=0.28, cell_rh=0.27)

        blurb = self._clean_text(texts.get("channels_summary", ""))
        if blurb:
            self._add_text(slide, TABLE_LEFT, bottom + 0.1, TABLE_W,
                           SLIDE_BOTTOM - bottom - 0.15, blurb, size=9)

    def _slide_channels_charts(self, prs, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, "Organic Channel Comparison: Detailed View")

        chart1 = str(self.charts_dir / "chart_google_direct.png")
        chart2 = str(self.charts_dir / "chart_channels_others.png")

        self._add_image(slide, chart1, TABLE_LEFT, CONTENT_TOP, (TABLE_W / 2) - 0.1, CONTENT_H)
        self._add_image(slide, chart2, TABLE_LEFT + (TABLE_W / 2) + 0.1, CONTENT_TOP, (TABLE_W / 2) - 0.1, CONTENT_H)

    def _slide_core_web_vitals(self, prs, data, month_label):
        slide = self._add_slide(prs)
        self._header_bar(slide, f"Core Web Vitals: {month_label}")

        cwv_data = data.get("core_web_vitals")
        if not cwv_data:
            self._add_rect(slide, TABLE_LEFT, CONTENT_TOP, TABLE_W, CONTENT_H, fill=C_GRAY_ALT)
            self._add_text(
                slide,
                TABLE_LEFT + 0.5,
                CONTENT_TOP + 1.5,
                TABLE_W - 1.0,
                1.5,
                "Core Web Vitals data unavailable.\n\nPlease verify that GOOGLE_API_KEY is configured in your .env file and has PageSpeed Insights API enabled in the Google Cloud Console.",
                size=12,
                color=C_DARK,
                align=PP_ALIGN.CENTER,
            )
            return

        metrics_keys = ["fcp", "lcp", "tbt", "cls", "si"]
        y_offset = CONTENT_TOP

        for key in metrics_keys:
            metric_data = cwv_data.get(key, {})
            status = metric_data.get("status", "unknown")
            text = metric_data.get("text", "")

            if not text:
                continue

            # Split header name/value line from description
            parts = text.split(". ", 1)
            header_line = parts[0] + "." if len(parts) > 1 else text
            desc_line = parts[1] if len(parts) > 1 else ""

            # Draw header line block depending on status
            if status == "good":
                fill_color = RGBColor(0xE2, 0xF0, 0xD9)
                text_color = RGBColor(0x38, 0x57, 0x23)
                self._add_rect(slide, TABLE_LEFT, y_offset, TABLE_W, 0.28, fill=fill_color)
                self._add_text(
                    slide,
                    TABLE_LEFT + 0.1,
                    y_offset + 0.02,
                    TABLE_W - 0.2,
                    0.24,
                    header_line,
                    size=10,
                    bold=True,
                    color=text_color,
                )
            elif status == "needs_improvement":
                fill_color = RGBColor(0xFF, 0xF2, 0xCC)
                text_color = RGBColor(0x7F, 0x60, 0x00)
                self._add_rect(slide, TABLE_LEFT, y_offset, TABLE_W, 0.28, fill=fill_color)
                self._add_text(
                    slide,
                    TABLE_LEFT + 0.1,
                    y_offset + 0.02,
                    TABLE_W - 0.2,
                    0.24,
                    header_line,
                    size=10,
                    bold=True,
                    color=text_color,
                )
            elif status == "poor":
                text_color = RGBColor(0xCC, 0x00, 0x00)
                self._add_text(
                    slide,
                    TABLE_LEFT,
                    y_offset + 0.02,
                    TABLE_W,
                    0.24,
                    header_line,
                    size=10,
                    bold=True,
                    color=text_color,
                )
            else:
                self._add_text(
                    slide,
                    TABLE_LEFT,
                    y_offset + 0.02,
                    TABLE_W,
                    0.24,
                    header_line,
                    size=10,
                    bold=True,
                    color=C_DARK,
                )

            # Draw description paragraph below with bolded numeric values
            if desc_line:
                self._add_desc_with_bolding(
                    slide,
                    TABLE_LEFT,
                    y_offset + 0.30,
                    TABLE_W,
                    0.50,
                    desc_line,
                    size=9,
                )

            y_offset += 0.92

    def _add_desc_with_bolding(self, slide, l, t, w, h, text, size=9.0, font=FONT_BODY):
        tb = slide.shapes.add_textbox(self._px(l), self._px(t), self._px(w), self._px(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = self._px(0)
        tf.margin_right = self._px(0)
        tf.margin_top = self._px(0)
        tf.margin_bottom = self._px(0)
        p = tf.paragraphs[0]

        # Regex to extract numeric values with units or standalone decimals
        pattern = re.compile(r"(\b\d+(?:\.\d+)?(?:ms|s)?\b)")
        parts = pattern.split(text)

        for part in parts:
            run = p.add_run()
            run.text = part
            run.font.size = Pt(size)
            run.font.name = font
            run.font.color.rgb = C_DARK
            if pattern.match(part):
                run.font.bold = True
            else:
                run.font.bold = False
        return tb
